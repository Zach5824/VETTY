"""Generate the Vetty presentation. Run: backend/.venv/bin/python scripts/create_presentation.py"""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "Vetty_Project_Presentation.pptx"
NAVY, BLUE, SKY, GOLD = "07517F", "1376B8", "248FCD", "F6C94C"
INK, MUTED, PALE, LINE, WHITE, GREEN, RED = "123047", "60768A", "F3F9FD", "D9E7F1", "FFFFFF", "29966D", "D95662"
prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]

def rgb(v): return RGBColor.from_string(v)
def box(s, x, y, w, h, fill, rounded=True, line=None):
    shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill); shape.line.color.rgb = rgb(line or fill)
    if rounded: shape.adjustments[0] = 0.1
    return shape
def write(s, value, x, y, w, h, size=16, fill=INK, bold=False, align=PP_ALIGN.LEFT, font="Aptos", italic=False):
    shape = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = shape.text_frame
    tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = align; run = p.add_run(); run.text = value
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = rgb(fill)
    return shape
def title(s, n, heading, sub="", dark=False):
    fg, faint = (WHITE, "D7EAF7") if dark else (INK, MUTED)
    box(s, .7, .48, .48, .48, GOLD); write(s, f"{n:02d}", .7, .64, .48, .13, 9, NAVY, True, PP_ALIGN.CENTER)
    write(s, heading, 1.38, .46, 11.1, .5, 27, fg, True, font="Aptos Display")
    if sub: write(s, sub, 1.4, 1.15, 11.1, .34, 12.5, faint)
def footer(s, page, dark=False):
    fill = "A9CCDF" if dark else MUTED
    write(s, "VETTY  •  PET CARE, SIMPLIFIED", .7, 7.08, 3.6, .16, 8, fill, True); write(s, f"{page:02d}", 12.1, 7.08, .5, .16, 8, fill, True, PP_ALIGN.RIGHT)
def panel(s, x, y, w, h, n, heading, body, accent=BLUE):
    box(s, x, y, w, h, WHITE, True, LINE); box(s, x+.25, y+.25, .46, .46, accent)
    write(s, str(n).zfill(2), x+.25, y+.41, .46, .12, 8, WHITE, True, PP_ALIGN.CENTER)
    write(s, heading, x+.25, y+.91, w-.5, .36, 14, INK, True); write(s, body, x+.25, y+1.42, w-.5, h-1.64, 10, MUTED)

# 1. Intro
s = prs.slides.add_slide(BLANK); box(s, 0, 0, 13.333, 7.5, NAVY, False); box(s, 9.45, -.85, 4.8, 4.8, SKY); box(s, 10.45, 4.1, 3.8, 3.8, GOLD)
box(s, .72, .72, 1.16, .34, GOLD); write(s, "PROJECT INTRO", .72, .82, 1.16, .11, 8, NAVY, True, PP_ALIGN.CENTER)
write(s, "Vetty", .72, 1.55, 6.2, .68, 42, WHITE, True, font="Aptos Display"); write(s, "Pet care, simplified.", .72, 2.4, 6.5, .42, 22, GOLD, True, font="Aptos Display", italic=True)
write(s, "A full-stack marketplace for Kenyan pet owners to shop for products, book veterinary services and pay safely in one mobile-first experience.", .72, 3.2, 6.55, .7, 15, "D7EAF7")
for x, label in [(.72, "Products + care"), (2.33, "M-Pesa + card"), (3.88, "Admin workspace")]: box(s, x, 4.4, 1.4, .33, "0E6598"); write(s, label, x, 4.5, 1.4, .12, 8.3, WHITE, True, PP_ALIGN.CENTER)
write(s, "Full-stack project • Kenya", .72, 6.7, 3.3, .18, 10, "A9CCDF", True)

# 2. Problem statement
s = prs.slides.add_slide(BLANK); title(s, 2, "Problem statement", "Pet-care shopping, health services and delivery are often coordinated in disconnected ways.")
panel(s, .7, 2.08, 3.82, 2.75, 1, "Scattered shopping", "Owners may visit several sellers to find food, toys and accessories, while also checking stock and delivery availability.", GOLD)
panel(s, 4.76, 2.08, 3.82, 2.75, 2, "Disconnected care", "Vet checkups, vaccinations and grooming bookings sit outside the product-buying journey, creating unnecessary friction.", SKY)
panel(s, 8.82, 2.08, 3.82, 2.75, 3, "Manual operations", "Staff need a secure central workspace for catalogue, inventory, delivery zones and fulfilment decisions.", GREEN)
write(s, "The opportunity: make pet care feel like one simple, trusted journey.", .7, 5.55, 11.7, .3, 15, NAVY, True); footer(s, 2)

# 3. Solution
s = prs.slides.add_slide(BLANK); box(s, 0, 0, 13.333, 7.5, PALE, False); title(s, 3, "Solution", "Vetty connects shopping, veterinary care, delivery and operations in one platform.")
for x, n, h, b, a in [(.7, 1, "Discover", "Browse products, services and delivery zones from one customer app.", BLUE), (3.82, 2, "Transact", "Build a cart, book appointments and use M-Pesa or card payment.", GOLD), (6.94, 3, "Track", "Keep orders, bookings and payment progress visible after checkout.", SKY), (10.06, 4, "Operate", "Give staff a protected admin workspace for day-to-day control.", GREEN)]: panel(s, x, 2.15, 2.58, 2.8, n, h, b, a)
write(s, "Customer-facing convenience + staff-facing control = a complete marketplace workflow.", .7, 5.65, 11.5, .3, 14, NAVY, True); footer(s, 3)

# 4. Development process
s = prs.slides.add_slide(BLANK); title(s, 4, "Development process", "A structured process kept the experience focused while the backend stayed secure and testable.")
steps = [("Discover", "Define the pet-owner and staff workflows."), ("Design", "Create a mobile-first customer and admin route map."), ("Build", "Develop React screens and Flask REST endpoints together."), ("Integrate", "Connect auth, persistence and payment-provider flows."), ("Verify", "Test role checks, ownership and payment scenarios.")]
for i, (h, b) in enumerate(steps):
    x=.7+i*2.5; box(s, x, 2.4, 2.1, 2.48, WHITE, True, LINE); write(s, str(i+1).zfill(2), x+.22, 2.7, .5, .16, 9, GOLD if i in (1,3) else BLUE, True); write(s, h, x+.22, 3.18, 1.65, .3, 14, INK, True); write(s, b, x+.22, 3.73, 1.64, .68, 9.5, MUTED)
    if i<4: write(s, "→", x+2.18, 3.43, .25, .25, 16, SKY, True, PP_ALIGN.CENTER)
footer(s, 4)

# 5. App features
s = prs.slides.add_slide(BLANK); title(s, 5, "App features", "Two experiences share the same source of truth: one for pet owners, one for operations.")
panel(s, .7, 2.08, 3.82, 3.2, 1, "Customer app", "Account registration and sign-in\nProduct discovery and cart\nService booking\nOrder and booking history", BLUE)
panel(s, 4.76, 2.08, 3.82, 3.2, 2, "Payments", "Stripe PaymentIntent flow\nM-Pesa STK Push flow\nVerified Stripe webhooks\nProtected M-Pesa callbacks", GOLD)
panel(s, 8.82, 2.08, 3.82, 3.2, 3, "Admin workspace", "Catalogue management\nInventory thresholds\nDelivery zones\nOrder and booking fulfilment", GREEN); footer(s, 5)

# 6. Techstack
s = prs.slides.add_slide(BLANK); box(s, 0, 0, 13.333, 7.5, NAVY, False); title(s, 6, "Tech stack", "Technologies selected for a fast mobile interface, maintainable API and production-ready payments.", True)
for x, label, tech, accent in [(.7, "FRONTEND", "React\nVite\nRedux Toolkit\nTailwind CSS", SKY), (3.82, "BACKEND", "Flask\nFlask-SQLAlchemy\nFlask-JWT-Extended\nREST API", GOLD), (6.94, "DATA", "SQLite for local work\nPostgreSQL for production\nSQLAlchemy ORM", GREEN), (10.06, "PAYMENTS", "Stripe.js + PaymentIntents\nSafaricom Daraja\nM-Pesa STK Push", SKY)]:
    box(s, x, 2.2, 2.58, 3.05, WHITE, True, LINE); write(s, label, x+.25, 2.55, 2, .18, 9, accent, True); write(s, tech, x+.25, 3.12, 2.02, 1.45, 13, INK, True)
footer(s, 6, True)

# 7. Feature implementations
s = prs.slides.add_slide(BLANK); title(s, 7, "Feature implementations", "The most important flows are enforced by the backend, not only represented in the interface.")
panel(s, .7, 2.08, 3.82, 3.12, 1, "Authentication & roles", "Passwords are hashed. JWT sessions identify the user. The API enforces customer ownership and administrator-only actions.", BLUE)
panel(s, 4.76, 2.08, 3.82, 3.12, 2, "Checkout & payments", "Orders are stored before payment. Stripe and M-Pesa payment status changes are validated through protected provider flows.", GOLD)
panel(s, 8.82, 2.08, 3.82, 3.12, 3, "Persistent admin tools", "Admin changes to products, services, thresholds, zones, orders and bookings are saved through REST endpoints and remain after refresh.", GREEN); footer(s, 7)

# 8. Challenges faced
s = prs.slides.add_slide(BLANK); title(s, 8, "Challenges faced", "Building a marketplace means resolving both user-experience and operational reliability concerns.")
panel(s, .7, 2.12, 3.82, 3.05, 1, "Secure access", "Role-based routes alone are not enough; authorization must be checked by every protected API action.", RED)
panel(s, 4.76, 2.12, 3.82, 3.05, 2, "Payment uncertainty", "Provider callbacks are asynchronous, so payment state must be verified before an order is treated as complete.", GOLD)
panel(s, 8.82, 2.12, 3.82, 3.05, 3, "Serverless persistence", "Local SQLite is not durable on serverless hosting; production requires a managed PostgreSQL database.", BLUE)
write(s, "Resolution: JWT ownership checks, verified provider callbacks and production PostgreSQL configuration.", .7, 5.68, 11.5, .3, 13.5, NAVY, True); footer(s, 8)

# 9. Video
s = prs.slides.add_slide(BLANK); box(s, 0, 0, 13.333, 7.5, PALE, False); title(s, 9, "Video demonstration", "Use this slide to play the recorded walkthrough or demonstrate the live application.")
box(s, 1.4, 1.95, 10.55, 4.18, NAVY, True); box(s, 5.88, 3.2, 1.6, 1.02, GOLD); write(s, "▶", 5.88, 3.41, 1.6, .35, 25, NAVY, True, PP_ALIGN.CENTER)
write(s, "Demo flow", 1.8, 2.37, 3, .28, 14, GOLD, True); write(s, "Customer: browse → book / checkout → pay → track\nAdmin: sign in → manage catalogue → approve fulfilment", 2, 4.77, 9.35, .62, 13, WHITE, True, PP_ALIGN.CENTER)
write(s, "Insert your screen-recorded video here, or use the live app during the presentation.", 1.8, 6.37, 9.8, .24, 10.5, MUTED, False, PP_ALIGN.CENTER); footer(s, 9)

# 10. Thanks statement
s = prs.slides.add_slide(BLANK); box(s, 0, 0, 13.333, 7.5, NAVY, False); box(s, 9.6, -.5, 4.25, 4.25, SKY); box(s, 10.75, 4.45, 3.6, 3.6, GOLD)
write(s, "Thank you", .75, 1.4, 7.3, .62, 34, WHITE, True, font="Aptos Display"); write(s, "Vetty brings pet products, trusted care and delivery together—while giving staff the tools to run the service well.", .75, 2.48, 6.7, .8, 17, "D7EAF7")
write(s, "Questions?", .75, 4.18, 2.7, .32, 17, GOLD, True); write(s, "Pet care, simplified.", .75, 6.5, 3.5, .3, 14, GOLD, True, font="Aptos Display", italic=True)
prs.core_properties.title = "Vetty — Project Presentation"; prs.core_properties.subject = "Full-stack pet-care marketplace"; prs.core_properties.author = "Vetty Project Team"
prs.save(OUT); print(f"Created {OUT}")
